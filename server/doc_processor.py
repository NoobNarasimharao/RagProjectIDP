import fitz  # PyMuPDF
from langchain_text_splitters import RecursiveCharacterTextSplitter
import re
import os
from pptx import Presentation
from docx import Document

def extract_text_from_pdf(pdf_path):
    """Extracts text from a PDF file and keeps track of page numbers."""
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF file not found at: {pdf_path}")

    doc = fitz.open(pdf_path)
    pages_content = []
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text")
        lines = text.split('\n')
        cleaned_lines = [line.strip() for line in lines if line.strip() and not re.match(r'^\d+$', line.strip())]
        pages_content.append((page_num, "\n".join(cleaned_lines)))
    doc.close()
    return pages_content

def extract_text_from_pptx(pptx_path):
    """Extracts text from a PPTX file. Each slide is treated as a page."""
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX file not found at: {pptx_path}")

    prs = Presentation(pptx_path)
    pages_content = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        pages_content.append((i, "\n".join(slide_text)))
    return pages_content

def extract_text_from_docx(docx_path):
    """Extracts text from a DOCX file. Treat as single page."""
    if not os.path.exists(docx_path):
        raise FileNotFoundError(f"DOCX file not found at: {docx_path}")

    doc = Document(docx_path)
    full_text = "\n".join([p.text for p in doc.paragraphs])
    return [(1, full_text)]

def extract_text_with_pages(file_path):
    """
    Generic extractor that handles PDF, PPTX, and DOCX.
    Returns a list of (page_number, text) tuples.
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    else:
        raise ValueError(f"Unsupported file extension: {ext}")

def chunk_text_with_metadata(pages_content, source_doc):
    """
    Chunks the extracted text and associates each chunk with its page number and index.
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=150,
        length_function=len,
        is_separator_regex=False,
    )

    all_chunks = []
    chunk_count = 0

    for page_num, text in pages_content:
        if not text.strip():
            continue
        page_chunks = text_splitter.split_text(text)
        for chunk in page_chunks:
            all_chunks.append({
                "text": chunk,
                "metadata": {
                    "source_pdf": source_doc, # Keep key as source_pdf for consistency with vector_store
                    "page_number": page_num,
                    "chunk_index": chunk_count
                }
            })
            chunk_count += 1

    return all_chunks
